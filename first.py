import os
from pptx import Presentation
from PyPDF2 import PdfReader
from langchain_core.documents import Document
from langchain.document_loaders import PyPDFLoader
from langchain_chroma import Chroma
from langchain.chains import RetrievalQA
from langchain.text_splitter import RecursiveCharacterTextSplitter
import shutil
import time
from langchain.llms.base import LLM
from langchain_groq import ChatGroq  
from langchain_mistralai.embeddings import MistralAIEmbeddings
from mistralai.client import MistralClient                                                                                                                                                   folder_path = "Data"
chroma_path = "chroma"
def generate_data():
    """Loads data from PPT/PPTX and PDF files, splits it, and saves it to ChromaDB."""
    ppt_documents = load_powerpoint_from_folder(folder_path)
    pdf_documents = load_pdf_from_folder(folder_path)
    all_documents = ppt_documents + pdf_documents
    chunks = split_text(all_documents)
    save_to_chroma(chunks)

def load_powerpoint_from_folder(folder_path):
    """Loads text content from all PPT/PPTX files in a folder into Langchain Documents."""
    documents = []
    for filename in os.listdir(folder_path):
        if filename.endswith((".ppt", ".pptx")):
            file_path = os.path.join(folder_path,filename)
            try:
                presentation = Presentation(file_path)
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text += run.text
                                text += "\n"  # Add a newline between paragraphs
                    text += "\n\n"  # Add a larger gap between slides
                metadata = {"source": file_path, "file_type": "PowerPoint"}
                documents.append(Document(page_content=text, metadata=metadata))
            except Exception as e:
                print(f"Error reading {filename}: {e}")
    return documents

def load_pdf_from_folder(folder_path):
    """Loads text content from all PDF files in a folder into Langchain Documents."""
    documents = []
    for filename in os.listdir(folder_path):
        if filename.endswith(".pdf"):
            file_path = os.path.join(folder_path, filename)
            try:
                loader = PyPDFLoader(file_path)
                pdf_documents = loader.load()
                for doc in pdf_documents:
                    doc.metadata["source"] = file_path
                    doc.metadata["file_type"] = "PDF"
                    documents.append(doc)
            except Exception as e:
                print(f"Error reading {filename}: {e}")
    return documents

def split_text(documents:list[Document]):
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = text_splitter.split_documents(documents)
        print(f"split {len(documents)} documents into {len(chunks)} chunks.")
        return chunks
def save_to_chroma(chunks: list[Document]):
    import warnings
    from langchain_mistralai import embeddings

    warnings.filterwarnings("ignore", message="Could not download mistral tokenizer from Huggingface")
    embeddings = MistralAIEmbeddings()
    db = Chroma.from_documents(chunks, embeddings, persist_directory=chroma_path)
    db = None
    print(f"saved {len(chunks)} chunks to {chroma_path}.")
def query_data(query):
    # Initialize the embedding model
    embeddings = MistralAIEmbeddings()

    # Load the ChromaDB
    db = Chroma(persist_directory=chroma_path, embedding_function=embeddings)

    # Create a retriever
    retriever = db.as_retriever(search_kwargs={'k': 4})
    relevant_documents = retriever.invoke(query)
    context = "\n\n".join([doc.page_content + f" (Source: {doc.metadata.get('file_type', 'unknown')})" for doc in relevant_documents])

    prompt = f"""You will try to answer the following question based on the provided documents.
Prioritize information found directly within the content of PowerPoint presentations or PDF documents.
If the answer is clearly and sufficiently present in the provided document content, use that information to answer.
Cite the source of your information by mentioning "(from PowerPoint)" or "(from PDF)".

If, after reviewing the document content, you cannot find a direct and complete answer, or if the information is insufficient, then you can use your general knowledge to provide a more comprehensive answer. In this case, do not cite a specific document.

Question: {query}

Document Content:
{context}

Answer: """

    try:
        groq_api_key = os.environ.get("GROQ_API_KEY")
        if not groq_api_key:
            raise ValueError("GROQ_API_KEY environment variable not set.")

        llm = ChatGroq(api_key=groq_api_key, model_name="llama3-70b-8192")
        answer = llm.invoke(prompt)  # This works for ChatGroq!
        print("\nAnswer for query: ")
        print(answer.content)

    except ImportError:
        print("\nLangchain Groq library not found. Please install it using 'pip install langchain-groq'.")
    except ValueError as e:
        print(f"\nError: {e}")
    except Exception as e:
        print(f"\nError during Groq API call: {e}")
    finally:
        del db

generate_data()                                                                                                                                                                                         def main(query):
    query_data(query)                                                                                                                                                                                       while True:
    query=input("Enter your Question(or type 'quit' to stop): ")
    if query.lower()=="quit":
        print("Hope this was helpful.")
        print("See you again soon.")
        break
    else:
        main(query)